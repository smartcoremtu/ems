"""
SmartCORE HEMS Presentation Generator
MTU / Interreg NWE SmartCORE
10-minute partner presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import pptx.util as util

# ── Brand colours ─────────────────────────────────────────────────────────────
MTU_BLUE      = RGBColor(0x00, 0x30, 0x87)   # MTU primary blue
MTU_GOLD      = RGBColor(0xF0, 0xAB, 0x00)   # MTU gold
SC_GREEN      = RGBColor(0x00, 0x96, 0x54)   # SmartCORE / Interreg green
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY    = RGBColor(0xF2, 0xF2, 0xF2)
DARK_GRAY     = RGBColor(0x33, 0x33, 0x33)
ACCENT_BLUE   = RGBColor(0x00, 0x70, 0xC0)

# ── Slide dimensions (widescreen 16:9) ────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # truly blank layout

# ── Helper utilities ──────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill_color, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(
        pptx.enum.shapes.MSO_SHAPE_TYPE.FREEFORM if False else 1,  # MSO_SHAPE.RECTANGLE
        l, t, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        if line_width:
            shape.line.width = line_width
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, l, t, w, h,
             font_size=18, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_para(tf, text, font_size=16, bold=False, color=DARK_GRAY,
             align=PP_ALIGN.LEFT, indent_level=0, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.level = indent_level
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def add_bullet_box(slide, bullets, l, t, w, h,
                   font_size=17, color=DARK_GRAY, heading=None):
    txBox = slide.shapes.add_textbox(l, t, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if heading:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = heading
        run.font.size = Pt(font_size + 2)
        run.font.bold = True
        run.font.color.rgb = MTU_BLUE
        first = False

    for bullet in bullets:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        indent = 0
        text = bullet
        if bullet.startswith("  –"):
            indent = 1
            text = bullet.strip()
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size - (2 if indent else 0))
        run.font.color.rgb = color
        p.level = indent
    return txBox

def header_bar(slide, title, subtitle=None, dark=True):
    """Full-width top header bar."""
    bar_h = Inches(1.45) if subtitle else Inches(1.1)
    add_rect(slide, 0, 0, W, bar_h, MTU_BLUE)
    add_text(slide, title,
             Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.85),
             font_size=32, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.4), Inches(0.88), Inches(12.5), Inches(0.45),
                 font_size=16, bold=False, color=MTU_GOLD, align=PP_ALIGN.LEFT)

def footer_bar(slide, left_text="MTU SmartCORE HEMS", right_text="Interreg NWE SmartCORE"):
    """Thin footer strip."""
    foot_h = Inches(0.3)
    foot_t = H - foot_h
    add_rect(slide, 0, foot_t, W, foot_h, MTU_BLUE)
    add_text(slide, left_text,
             Inches(0.2), foot_t, Inches(6), foot_h,
             font_size=10, color=WHITE, align=PP_ALIGN.LEFT)
    add_text(slide, right_text,
             Inches(7), foot_t, Inches(6), foot_h,
             font_size=10, color=MTU_GOLD, align=PP_ALIGN.RIGHT)

def accent_line(slide, t=None):
    """Thin gold accent line below header."""
    if t is None:
        t = Inches(1.45)
    line = slide.shapes.add_shape(1, 0, t, W, Inches(0.04))
    line.fill.solid()
    line.fill.fore_color.rgb = MTU_GOLD
    line.line.fill.background()

def info_box(slide, title, body_lines, l, t, w, h, accent=SC_GREEN):
    """Coloured info card."""
    add_rect(slide, l, t, w, h, LIGHT_GRAY, line_color=accent, line_width=Pt(2))
    add_rect(slide, l, t, w, Inches(0.38), accent)  # top strip
    add_text(slide, title, l + Inches(0.12), t + Inches(0.02),
             w - Inches(0.2), Inches(0.36),
             font_size=13, bold=True, color=WHITE)
    txBox = slide.shapes.add_textbox(l + Inches(0.12), t + Inches(0.42),
                                     w - Inches(0.2), h - Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in body_lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(13)
        run.font.color.rgb = DARK_GRAY

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
# Full blue background
add_rect(s1, 0, 0, W, H, MTU_BLUE)
# Gold diagonal accent strip
shape = s1.shapes.add_shape(1, Inches(0), Inches(5.2), W, Inches(0.06))
shape.fill.solid(); shape.fill.fore_color.rgb = MTU_GOLD
shape.line.fill.background()

add_text(s1, "SmartCORE HEMS Box",
         Inches(0.6), Inches(1.4), Inches(12), Inches(1.6),
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "Home Energy Management System — Research Development",
         Inches(0.6), Inches(2.8), Inches(12), Inches(0.7),
         font_size=22, bold=False, color=MTU_GOLD, align=PP_ALIGN.LEFT)
add_text(s1, "Munster Technological University  |  Interreg NWE SmartCORE",
         Inches(0.6), Inches(3.55), Inches(12), Inches(0.5),
         font_size=17, bold=False, color=WHITE, align=PP_ALIGN.LEFT)
add_text(s1, "[Presenter Name]  ·  [Date]  ·  SmartCORE Partner Meeting",
         Inches(0.6), Inches(4.1), Inches(12), Inches(0.4),
         font_size=14, italic=True, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)

# Logo placeholder boxes
add_rect(s1, Inches(0.6), Inches(5.5), Inches(2.6), Inches(1.3),
         RGBColor(0x00, 0x20, 0x60), line_color=MTU_GOLD, line_width=Pt(1.5))
add_text(s1, "[ MTU Logo ]", Inches(0.6), Inches(5.7), Inches(2.6), Inches(0.9),
         font_size=14, color=MTU_GOLD, align=PP_ALIGN.CENTER)

add_rect(s1, Inches(3.5), Inches(5.5), Inches(3.2), Inches(1.3),
         RGBColor(0x00, 0x20, 0x60), line_color=SC_GREEN, line_width=Pt(1.5))
add_text(s1, "[ SmartCORE NWE / Interreg Logo ]",
         Inches(3.5), Inches(5.7), Inches(3.2), Inches(0.9),
         font_size=11, color=SC_GREEN, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — SMARTCORE NWE CONTEXT
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, W, H, WHITE)
header_bar(s2, "SmartCORE NWE — Project Context",
           "Interreg North-West Europe  |  MTU's Role")
accent_line(s2)
footer_bar(s2)

bullets_left = [
    "• Interreg NWE transnational research project",
    "• Partners across Ireland, UK, Belgium, France, Netherlands",
    "• Focus: Community-scale smart energy management",
    "• Objective: Reduce household energy costs & carbon footprint",
    "",
    "• MTU leads the Ireland pilot workstream",
    "• Responsible for: HEMS hardware + software stack development",
    "• Target: Deploy HEMS units in Irish pilot homes",
]
add_bullet_box(s2, bullets_left, Inches(0.4), Inches(1.6), Inches(6.0), Inches(5.4),
               font_size=17)

info_box(s2, "MTU Research Role",
         ["Design & build open-source HEMS gateway",
          "Containerised, remotely manageable",
          "Bridges Zigbee sensors → cloud analytics",
          "Feeds SmartCORE data platform"],
         Inches(6.7), Inches(1.7), Inches(6.2), Inches(2.3), accent=MTU_GOLD)

info_box(s2, "Key Metrics (Target)",
         ["≥ 20 pilot homes in Ireland",
          "Real-time energy monitoring",
          "15-min resolution data to EU platform",
          "OTA updates — no on-site visits"],
         Inches(6.7), Inches(4.2), Inches(6.2), Inches(2.3), accent=SC_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, W, H, WHITE)
header_bar(s3, "The Challenge", "Why do we need a bespoke HEMS gateway?")
accent_line(s3)
footer_bar(s3)

problems = [
    ("Fragmented Sensor Ecosystem",
     ["Zigbee, Z-Wave, BLE — no common protocol",
      "Utility meters rarely IoT-ready",
      "Off-the-shelf hubs lack research-grade data access"]),
    ("Data Sovereignty & Privacy",
     ["Cloud-dependent products = no local control",
      "GDPR concerns with home energy data",
      "Partners need raw data, not vendor-filtered APIs"]),
    ("Remote Manageability at Scale",
     ["30+ homes across Ireland — no on-site engineers",
      "Firmware updates must be zero-touch",
      "Self-healing critical — no single point of failure"]),
    ("Open & Extensible Research Platform",
     ["Off-the-shelf systems are black boxes",
      "We need to add sensors, change algorithms, export data",
      "Must integrate with SmartCORE EU data platform"]),
]

positions = [
    (Inches(0.3),  Inches(1.55)),
    (Inches(6.7),  Inches(1.55)),
    (Inches(0.3),  Inches(4.3)),
    (Inches(6.7),  Inches(4.3)),
]
for (title, lines), (l, t) in zip(problems, positions):
    info_box(s3, title, lines, l, t, Inches(6.1), Inches(2.5), accent=MTU_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — WHAT IS THE HEMS BOX?
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, W, H, WHITE)
header_bar(s4, "The HEMS Box — What Is It?",
           "An open, containerised edge gateway for residential energy monitoring")
accent_line(s4)
footer_bar(s4)

# Hardware column
add_rect(s4, Inches(0.3), Inches(1.6), Inches(4.0), Inches(5.5), LIGHT_GRAY,
         line_color=MTU_BLUE, line_width=Pt(2))
add_rect(s4, Inches(0.3), Inches(1.6), Inches(4.0), Inches(0.42), MTU_BLUE)
add_text(s4, "Hardware", Inches(0.42), Inches(1.62), Inches(3.8), Inches(0.4),
         font_size=14, bold=True, color=WHITE)

hw = [
    "Raspberry Pi 5  (arm64 / aarch64)",
    "  – 4 GB RAM, Quad-core Cortex-A76",
    "  – Runs BalenaOS (Yocto Linux)",
    "",
    "Sonoff ZBDongle-E (USB)",
    "  – EZSP Zigbee coordinator",
    "  – IEEE 802.15.4 @ 2.4 GHz",
    "",
    "Frient SMSZB-120",
    "  – Smart meter interface sensor",
    "  – Blink counter → power & energy",
    "",
    "[ Photo placeholder: Pi 5 + USB dongle ]",
]
add_bullet_box(s4, hw, Inches(0.42), Inches(2.1), Inches(3.7), Inches(4.7),
               font_size=14)

# Software column
add_rect(s4, Inches(4.6), Inches(1.6), Inches(4.0), Inches(5.5), LIGHT_GRAY,
         line_color=SC_GREEN, line_width=Pt(2))
add_rect(s4, Inches(4.6), Inches(1.6), Inches(4.0), Inches(0.42), SC_GREEN)
add_text(s4, "Software Stack", Inches(4.72), Inches(1.62), Inches(3.8), Inches(0.4),
         font_size=14, bold=True, color=WHITE)

sw = [
    "BalenaOS — container OS + OTA",
    "Docker Compose — 8 services",
    "",
    "Home Assistant — data hub & UI",
    "InfluxDB 2 — time-series database",
    "Nginx — reverse proxy & auth",
    "Mosquitto — MQTT broker",
    "hass-configurator — config UI",
    "system-manager — watchdog",
    "led-status — GPIO health LEDs",
    "",
    "Bridge network 172.18.4.0/24",
    "Isolated, static IPs per service",
]
add_bullet_box(s4, sw, Inches(4.72), Inches(2.1), Inches(3.7), Inches(4.7),
               font_size=14)

# Principle column
add_rect(s4, Inches(8.9), Inches(1.6), Inches(4.1), Inches(5.5), LIGHT_GRAY,
         line_color=MTU_GOLD, line_width=Pt(2))
add_rect(s4, Inches(8.9), Inches(1.6), Inches(4.1), Inches(0.42), MTU_GOLD)
add_text(s4, "Design Principles", Inches(9.02), Inches(1.62), Inches(3.9), Inches(0.4),
         font_size=14, bold=True, color=DARK_GRAY)

pr = [
    "Open-source throughout",
    "No cloud dependency",
    "  (works fully offline)",
    "",
    "Self-healing & self-updating",
    "  via BalenaOS Supervisor",
    "",
    "Extensible — add sensors,",
    "  services, or integrations",
    "  without hardware changes",
    "",
    "Research-grade data access",
    "  — full InfluxDB query API",
]
add_bullet_box(s4, pr, Inches(9.02), Inches(2.1), Inches(3.8), Inches(4.7),
               font_size=14)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, W, H, WHITE)
header_bar(s5, "System Architecture",
           "8 containerised services on a private Docker bridge network")
accent_line(s5)
footer_bar(s5)

# ── Architecture diagram (drawn as shapes) ────────────────────────────────────
# Physical layer label
add_text(s5, "Raspberry Pi 5  (BalenaOS Host)",
         Inches(0.2), Inches(1.55), Inches(12.9), Inches(0.35),
         font_size=13, bold=True, color=MTU_BLUE, align=PP_ALIGN.CENTER)

# Outer Pi box
pi_box = s5.shapes.add_shape(1, Inches(0.2), Inches(1.5), Inches(12.9), Inches(5.5))
pi_box.fill.solid(); pi_box.fill.fore_color.rgb = RGBColor(0xE8, 0xF0, 0xFF)
pi_box.line.color.rgb = MTU_BLUE; pi_box.line.width = Pt(2)

# Docker bridge label
add_text(s5, "Docker Bridge  172.18.4.0/24",
         Inches(0.4), Inches(1.9), Inches(8), Inches(0.3),
         font_size=11, italic=True, color=MTU_BLUE)

# Service boxes — row 1
services_r1 = [
    ("Home\nAssistant", "172.18.4.2", ":8123", MTU_BLUE),
    ("InfluxDB 2", "172.18.4.3", ":8086", SC_GREEN),
    ("Nginx\nProxy", "172.18.4.4", ":80", ACCENT_BLUE),
    ("Mosquitto\nMQTT", "172.18.4.7", ":1883", RGBColor(0x8B, 0x00, 0x8B)),
]
box_w = Inches(2.4); box_h = Inches(1.3)
row1_t = Inches(2.25)
gap = Inches(0.18)
start_l = Inches(0.4)
for i, (name, ip, port, col) in enumerate(services_r1):
    l = start_l + i * (box_w + gap)
    add_rect(s5, l, row1_t, box_w, box_h, col)
    add_text(s5, name, l + Inches(0.08), row1_t + Inches(0.05),
             box_w - Inches(0.15), Inches(0.7),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s5, f"{ip}\n{port}", l + Inches(0.08), row1_t + Inches(0.72),
             box_w - Inches(0.15), Inches(0.55),
             font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

# Service boxes — row 2
services_r2 = [
    ("hass-\nconfigurator", "172.18.4.6", ":3218", RGBColor(0xC0, 0x60, 0x00)),
    ("led-status", "172.18.4.9", "GPIO 17/27/22", RGBColor(0x55, 0x77, 0x00)),
    ("system-\nmanager", "host network", "Supervisor API", RGBColor(0x44, 0x44, 0x88)),
    ("wifi-\nconnect", "host network", "Captive portal", RGBColor(0x44, 0x88, 0x88)),
]
row2_t = Inches(3.75)
for i, (name, ip, port, col) in enumerate(services_r2):
    l = start_l + i * (box_w + gap)
    add_rect(s5, l, row2_t, box_w, box_h, col)
    add_text(s5, name, l + Inches(0.08), row2_t + Inches(0.05),
             box_w - Inches(0.15), Inches(0.7),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s5, f"{ip}\n{port}", l + Inches(0.08), row2_t + Inches(0.72),
             box_w - Inches(0.15), Inches(0.55),
             font_size=11, color=RGBColor(0xDD, 0xEE, 0xFF), align=PP_ALIGN.CENTER)

# External / host items on right
add_rect(s5, Inches(10.5), Inches(2.25), Inches(2.5), Inches(0.9),
         RGBColor(0xAA, 0x44, 0x00))
add_text(s5, "BalenaOS\nSupervisor", Inches(10.5), Inches(2.25), Inches(2.5), Inches(0.9),
         font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s5, Inches(10.5), Inches(3.35), Inches(2.5), Inches(0.9),
         RGBColor(0x55, 0x55, 0x55))
add_text(s5, "Balena Cloud\n(OTA updates)", Inches(10.5), Inches(3.35), Inches(2.5), Inches(0.9),
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# USB Dongle + Sensor (bottom)
add_rect(s5, Inches(0.4), Inches(5.25), Inches(3.2), Inches(0.65),
         RGBColor(0x33, 0x66, 0x00))
add_text(s5, "Sonoff ZBDongle-E (USB / EZSP Zigbee coordinator)",
         Inches(0.4), Inches(5.25), Inches(3.2), Inches(0.65),
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s5, Inches(4.0), Inches(5.25), Inches(3.0), Inches(0.65),
         RGBColor(0x33, 0x66, 0x00))
add_text(s5, "Frient SMSZB-120 (Zigbee sensor)",
         Inches(4.0), Inches(5.25), Inches(3.0), Inches(0.65),
         font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Arrow label
add_text(s5, "← Zigbee 802.15.4 →",
         Inches(3.25), Inches(5.32), Inches(0.75), Inches(0.5),
         font_size=10, italic=True, color=SC_GREEN)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — ENERGY DATA FLOW
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
add_rect(s6, 0, 0, W, H, WHITE)
header_bar(s6, "Energy Data Flow",
           "From physical smart meter pulse to database — end-to-end")
accent_line(s6)
footer_bar(s6)

# Flow boxes
flow_steps = [
    ("Smart Meter", "IR pulse output\n(blink counter)", RGBColor(0x33, 0x66, 0x00)),
    ("Frient\nSMSZB-120", "Counts pulses\n→ kWh / W", RGBColor(0x00, 0x77, 0x44)),
    ("Sonoff\nZBDongle-E", "USB Zigbee\ncoordinator\n(EZSP)", RGBColor(0x00, 0x60, 0x80)),
    ("ZHA\n(in HA)", "Zigbee Home\nAutomation\nintegration", MTU_BLUE),
    ("HA State\nMachine", "Entity states\nupdated every\n60 s", ACCENT_BLUE),
    ("InfluxDB\nIntegration", "HTTP POST\nLine Protocol\n→ InfluxDB", SC_GREEN),
    ("InfluxDB 2", "Time-series DB\nbucket:\nhome_assistant", RGBColor(0x00, 0x80, 0x40)),
]

box_w = Inches(1.62); box_h = Inches(2.3)
arrow_w = Inches(0.35)
total = len(flow_steps) * box_w + (len(flow_steps) - 1) * arrow_w
start_l = (W - total) / 2
row_t = Inches(2.1)

for i, (title, desc, col) in enumerate(flow_steps):
    l = start_l + i * (box_w + arrow_w)
    # Box
    add_rect(s6, l, row_t, box_w, box_h, col)
    add_text(s6, title, l + Inches(0.05), row_t + Inches(0.08),
             box_w - Inches(0.1), Inches(0.65),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s6, desc, l + Inches(0.05), row_t + Inches(0.75),
             box_w - Inches(0.1), Inches(1.45),
             font_size=11, color=RGBColor(0xDD, 0xFF, 0xDD), align=PP_ALIGN.CENTER)
    # Arrow (not after last)
    if i < len(flow_steps) - 1:
        arr_l = l + box_w
        add_text(s6, "→", arr_l, row_t + Inches(0.9),
                 arrow_w, Inches(0.5),
                 font_size=20, bold=True, color=MTU_GOLD, align=PP_ALIGN.CENTER)

# Protocol annotations beneath
protocols = [
    (start_l, "Zigbee 802.15.4", SC_GREEN),
    (start_l + 2 * (box_w + arrow_w), "EZSP/USB serial", ACCENT_BLUE),
    (start_l + 3 * (box_w + arrow_w), "HA internal bus", MTU_BLUE),
    (start_l + 5 * (box_w + arrow_w), "HTTP /api/v2/write\n(Line Protocol)", SC_GREEN),
]
for l, label, col in protocols:
    add_text(s6, label, l, row_t + box_h + Inches(0.15),
             box_w + arrow_w, Inches(0.5),
             font_size=11, italic=True, color=col, align=PP_ALIGN.CENTER)

# Key note
add_rect(s6, Inches(0.3), Inches(5.6), Inches(12.7), Inches(0.75),
         RGBColor(0xFF, 0xF8, 0xE0), line_color=MTU_GOLD, line_width=Pt(1.5))
add_text(s6,
         "Note: MQTT broker is present but idle in the current ZHA deployment. "
         "Zigbee2MQTT is an alternative architecture that routes all data through MQTT (see Slide 10).",
         Inches(0.45), Inches(5.65), Inches(12.4), Inches(0.65),
         font_size=12, italic=True, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — KEY SERVICES DEEP DIVE
# ══════════════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(BLANK)
add_rect(s7, 0, 0, W, H, WHITE)
header_bar(s7, "Key Services — What Each One Does")
accent_line(s7, Inches(1.1))
footer_bar(s7)

services_detail = [
    ("Home Assistant", MTU_BLUE,
     ["• Central data hub & user dashboard",
      "• Manages Zigbee devices via ZHA integration",
      "• Syncs all sensor states → InfluxDB automatically",
      "• Automation engine (backups, alerts)",
      "• Trusted proxy via Nginx"]),
    ("InfluxDB 2", SC_GREEN,
     ["• Time-series database for all energy data",
      "• Bucket: home_assistant  |  Org: hems",
      "• Queried via Flux language / REST API",
      "• Retains full history for research analysis",
      "• UI accessible at /influx/ via Nginx"]),
    ("Nginx Reverse Proxy", ACCENT_BLUE,
     ["• Single entry point — port 80 on host",
      "• Basic auth (CONFIG_USER / CONFIG_PASSWORD)",
      "• Routes /  →  HA, /influx/  →  InfluxDB",
      "• WebSocket upgrade headers for HA live UI",
      "• Sub-filter rewrites InfluxDB JS asset paths"]),
    ("system-manager\n(Watchdog)", RGBColor(0x44, 0x44, 0x88),
     ["• Runs every 30 s on host network",
      "• Restarts HA on new BalenaOS release",
      "• Reboots device if internet down > 30 min",
      "• Uses Balena Supervisor REST API",
      "• iptables rules: Docker ↔ WiFi repeater"]),
    ("Mosquitto MQTT", RGBColor(0x8B, 0x00, 0x8B),
     ["• Message broker (port 1883)",
      "• Ready for Zigbee2MQTT alternative arch",
      "• Retained messages → sensor state on restart",
      "• Persistence to mosquitto.db volume",
      "• Currently idle (ZHA deployment)"]),
    ("led-status", RGBColor(0x55, 0x77, 0x00),
     ["• GPIO 17 → green LED (HA reachable?)",
      "• GPIO 27 → blue LED (internet alive?)",
      "• GPIO 22 → red LED (HA errors in log?)",
      "• Pi 5 lgpio patch (chip=0)",
      "• Polls every 2 seconds"]),
]

cols = 3; rows = 2
cell_w = Inches(4.2); cell_h = Inches(2.55)
col_gap = Inches(0.15); row_gap = Inches(0.15)
start_l2 = Inches(0.3); start_t2 = Inches(1.2)
for idx, (name, col, lines) in enumerate(services_detail):
    ci = idx % cols; ri = idx // cols
    l = start_l2 + ci * (cell_w + col_gap)
    t = start_t2 + ri * (cell_h + row_gap)
    info_box(s7, name, lines, l, t, cell_w, cell_h, accent=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SELF-HEALING & REMOTE MANAGEMENT
# ══════════════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(BLANK)
add_rect(s8, 0, 0, W, H, WHITE)
header_bar(s8, "Self-Healing & Remote Management",
           "Zero-touch operation across distributed pilot deployments")
accent_line(s8)
footer_bar(s8)

# BalenaOS box
info_box(s8, "BalenaOS — OTA Updates",
         ["Yocto-based container OS with Supervisor daemon",
          "Push new Docker images from Balena Cloud → device auto-pulls",
          "No SSH access required — fully managed via Balena dashboard",
          "Balena App ID / Release ID tracked for graceful HA restart"],
         Inches(0.3), Inches(1.6), Inches(6.1), Inches(2.35), accent=MTU_BLUE)

info_box(s8, "system-manager Watchdog (30 s loop)",
         ["1. Stops wifi-repeater after 10 min (bandwidth management)",
          "2. Detects new Balena release → graceful HA restart (2 min grace)",
          "3. Pings google.com — reboots device if down > 30 min",
          "4. Sets iptables rules on startup (Docker ↔ WiFi subnets)"],
         Inches(6.6), Inches(1.6), Inches(6.4), Inches(2.35), accent=RGBColor(0x44, 0x44, 0x88))

info_box(s8, "LED Status Indicators (Physical Feedback)",
         ["🟢 GPIO 17 — Home Assistant reachable (internal ping)",
          "🔵 GPIO 27 — Internet connectivity (8.8.8.8 ping)",
          "🔴 GPIO 22 — Error state (scans HA log for ERROR keyword)",
          "Allows non-technical householders to see system health at a glance"],
         Inches(0.3), Inches(4.1), Inches(6.1), Inches(2.35), accent=SC_GREEN)

info_box(s8, "WiFi Onboarding — wifi-connect",
         ["On first boot (or lost WiFi): launches captive portal",
          "Householder connects phone → selects home WiFi → enters password",
          "NetworkManager configures WiFi via D-Bus on host",
          "No LAN cable or technical knowledge required"],
         Inches(6.6), Inches(4.1), Inches(6.4), Inches(2.35), accent=MTU_GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — DASHBOARD & DATA
# ══════════════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(BLANK)
add_rect(s9, 0, 0, W, H, WHITE)
header_bar(s9, "Energy Monitoring Dashboard",
           "Real-time visibility into household energy consumption")
accent_line(s9)
footer_bar(s9)

# Large screenshot placeholder
add_rect(s9, Inches(0.3), Inches(1.55), Inches(8.5), Inches(5.45),
         LIGHT_GRAY, line_color=MTU_BLUE, line_width=Pt(2))
add_text(s9, "[ Screenshot Placeholder ]\n\nHome Assistant Energy Dashboard\nInfluxDB Grafana-style charts\nReal-time power (W) & cumulative energy (kWh)",
         Inches(0.3), Inches(2.8), Inches(8.5), Inches(3.0),
         font_size=18, color=RGBColor(0x99, 0x99, 0x99), align=PP_ALIGN.CENTER)

# Right side metrics
info_box(s9, "Data Captured",
         ["Power (W) — real-time consumption",
          "Energy (kWh) — cumulative total",
          "Link quality — Zigbee signal strength",
          "Timestamps at 60 s resolution",
          "Full history in InfluxDB"],
         Inches(9.1), Inches(1.55), Inches(3.9), Inches(2.4), accent=SC_GREEN)

info_box(s9, "Access Points",
         ["Local LAN: http://<pi-ip>/",
          "InfluxDB UI: http://<pi-ip>/influx/",
          "Config UI: http://<pi-ip>/configurator/",
          "All behind Nginx basic auth",
          "WebSocket live updates supported"],
         Inches(9.1), Inches(4.1), Inches(3.9), Inches(2.4), accent=MTU_BLUE)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — EXTENSIBILITY & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(BLANK)
add_rect(s10, 0, 0, W, H, WHITE)
header_bar(s10, "Extensibility & Future Work",
           "The HEMS box is designed to grow with research requirements")
accent_line(s10)
footer_bar(s10)

ext_cards = [
    ("Alternative: Zigbee2MQTT", ACCENT_BLUE,
     ["Swap ZHA for Zigbee2MQTT container",
      "All data routed through MQTT broker",
      "Enables multi-subscriber architecture",
      "→ Edge ML, alerting, custom exporters",
      "MQTT Discovery → HA auto-creates entities"]),
    ("Mobile Access — Tailscale VPN", MTU_BLUE,
     ["Add tailscale container (WireGuard VPN)",
      "Researcher/admin accesses dashboard remotely",
      "Peer-to-peer — no cloud server needed",
      "Stable Tailscale IP survives dynamic WAN IP",
      "TS_AUTH_KEY via Balena env var (never in repo)"]),
    ("Mobile Access — Cloudflare Tunnel", SC_GREEN,
     ["Add cloudflared container",
      "4 outbound HTTPS connections to Cloudflare",
      "Public HTTPS URL: hems.yourdomain.com",
      "No port forwarding or router changes",
      "Cloudflare Access for pre-auth layer (free)"]),
    ("Additional Sensors", MTU_GOLD,
     ["Any Zigbee sensor pairs via ZHA/Z2M",
      "Temperature, humidity, occupancy, gas",
      "EV charger monitoring (eesmart-d2l path)",
      "Clamp-on CT sensors via ESPHome/WiFi",
      "Solar PV integration (planned)"]),
    ("Data Export & Integration", RGBColor(0x44, 0x44, 0x88),
     ["InfluxDB REST API → SmartCORE EU platform",
      "Export to CSV / Parquet for research",
      "Grafana compatible (Flux queries)",
      "MQTT → Node-RED for custom pipelines",
      "REST hook → external ML inference"]),
    ("Known Gaps to Resolve", RGBColor(0xAA, 0x33, 0x00),
     ["eesmart-d2l: TCP proxy in Nginx but no container",
      "No TLS on Nginx (plain HTTP on LAN only)",
      "MQTT allow_anonymous (harden for production)",
      "template_sensors/ not in git (on-device only)",
      "Pilot home deployment tooling — in progress"]),
]

cell_w = Inches(4.1); cell_h = Inches(2.4)
for idx, (name, col, lines) in enumerate(ext_cards):
    ci = idx % 3; ri = idx // 3
    l = Inches(0.3) + ci * (cell_w + Inches(0.2))
    t = Inches(1.2) + ri * (cell_h + Inches(0.15))
    info_box(s10, name, lines, l, t, cell_w, cell_h, accent=col)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — MTU RESEARCH CONTRIBUTION
# ══════════════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(BLANK)
add_rect(s11, 0, 0, W, H, WHITE)
header_bar(s11, "MTU's Research Contribution",
           "What we bring to SmartCORE NWE")
accent_line(s11)
footer_bar(s11)

# Left: what we built
add_bullet_box(s11,
    ["What MTU Has Built",
     "",
     "• Full open-source HEMS software stack",
     "  – Containerised, reproducible, version-controlled",
     "  – Deployable from a single docker-compose.yml",
     "",
     "• Raspberry Pi 5 validated deployment",
     "  – Pi 5 GPIO lgpio patch for led-status",
     "  – aarch64/arm64 Docker image selection",
     "",
     "• Self-healing watchdog architecture",
     "  – Balena Supervisor API integration",
     "  – Graceful HA restart on new firmware",
     "",
     "• Technical documentation & architecture reference",
     "  – ARCHITECTURE.md, COMPONENTS.md",
     "  – MQTT & Mobile Access reference",
     "  – PlantUML system diagrams",
    ],
    Inches(0.4), Inches(1.25), Inches(5.8), Inches(5.9),
    font_size=15, heading=None)

# Right: impact
info_box(s11, "Impact for SmartCORE Partners",
         ["Any partner can clone repo & deploy",
          "No proprietary hardware lock-in",
          "OTA updates mean pilot homes update silently",
          "Full data access — no vendor API restrictions",
          "Adapts to local smart meter standards"],
         Inches(6.4), Inches(1.3), Inches(6.5), Inches(2.1), accent=SC_GREEN)

info_box(s11, "Research Outputs (Planned)",
         ["Peer-reviewed publication on HEMS architecture",
          "Open dataset from pilot homes (anonymised)",
          "Comparison: ZHA vs Zigbee2MQTT performance",
          "Energy baseline per household profile",
          "Integration template for EU SmartCORE platform"],
         Inches(6.4), Inches(3.6), Inches(6.5), Inches(2.1), accent=MTU_BLUE)

info_box(s11, "Technology Readiness",
         ["TRL 4 → 5 transition (lab → pilot validated)",
          "Core stack functional on Pi 5 hardware",
          "Sensor integration confirmed end-to-end",
          "Pilot deployment tooling in progress"],
         Inches(6.4), Inches(5.85), Inches(6.5), Inches(1.2), accent=MTU_GOLD)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — STATUS & ROADMAP
# ══════════════════════════════════════════════════════════════════════════════
s12 = prs.slides.add_slide(BLANK)
add_rect(s12, 0, 0, W, H, WHITE)
header_bar(s12, "Current Status & Roadmap",
           "Early prototype → pilot deployment")
accent_line(s12)
footer_bar(s12)

# Status bar
statuses = [
    ("Stack Running", "All 8 containers\noperational on Pi 5", SC_GREEN, "✓"),
    ("Sensor Data", "Frient → InfluxDB\nend-to-end confirmed", SC_GREEN, "✓"),
    ("OTA Updates", "BalenaOS push\nvalidated in lab", SC_GREEN, "✓"),
    ("Self-Healing", "Watchdog tested\n(reboot, HA restart)", SC_GREEN, "✓"),
    ("Mobile Access", "Architecture designed,\nimplementation pending", MTU_GOLD, "~"),
    ("Pilot Homes", "Site survey in progress,\ndeployment Q3 2026", MTU_BLUE, "→"),
]
sw2 = Inches(2.0); sh2 = Inches(1.4); sg = Inches(0.12)
st = Inches(1.6)
for i, (title, desc, col, icon) in enumerate(statuses):
    l = Inches(0.3) + i * (sw2 + sg)
    add_rect(s12, l, st, sw2, sh2, col)
    add_text(s12, icon + " " + title, l + Inches(0.06), st + Inches(0.05),
             sw2 - Inches(0.1), Inches(0.52),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s12, desc, l + Inches(0.06), st + Inches(0.58),
             sw2 - Inches(0.1), Inches(0.75),
             font_size=11, color=RGBColor(0xDD, 0xFF, 0xDD), align=PP_ALIGN.CENTER)

# Roadmap
add_text(s12, "Roadmap", Inches(0.4), Inches(3.2), Inches(12), Inches(0.4),
         font_size=18, bold=True, color=MTU_BLUE)

roadmap = [
    ("Q2 2026", "Pilot preparation",
     "Finalise deployment tooling\nMobile access (Tailscale / Cloudflare)\nHarden MQTT security"),
    ("Q3 2026", "Pilot deployment",
     "Install HEMS in first 5–10 Irish homes\nLive data flowing to SmartCORE platform\nOn-site WiFi onboarding validation"),
    ("Q4 2026", "Scale & analyse",
     "Expand to 20 homes\nBaseline energy profiles\nPeer-reviewed publication draft"),
    ("2027", "Research outputs",
     "Full dataset analysis\nHEMS v2 design (lessons learned)\nResults shared with all NWE partners"),
]
rm_w = Inches(3.1); rm_h = Inches(2.7); rm_gap = Inches(0.15)
rm_t = Inches(3.65)
for i, (period, phase, detail) in enumerate(roadmap):
    l = Inches(0.3) + i * (rm_w + rm_gap)
    col = [MTU_GOLD, SC_GREEN, ACCENT_BLUE, MTU_BLUE][i]
    add_rect(s12, l, rm_t, rm_w, rm_h, LIGHT_GRAY, line_color=col, line_width=Pt(2))
    add_rect(s12, l, rm_t, rm_w, Inches(0.38), col)
    add_text(s12, period, l + Inches(0.08), rm_t + Inches(0.02),
             rm_w - Inches(0.15), Inches(0.36),
             font_size=14, bold=True, color=WHITE)
    add_text(s12, phase, l + Inches(0.08), rm_t + Inches(0.44),
             rm_w - Inches(0.15), Inches(0.42),
             font_size=13, bold=True, color=col)
    add_text(s12, detail, l + Inches(0.08), rm_t + Inches(0.9),
             rm_w - Inches(0.15), Inches(1.7),
             font_size=12, color=DARK_GRAY)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — THANK YOU / Q&A
# ══════════════════════════════════════════════════════════════════════════════
s13 = prs.slides.add_slide(BLANK)
add_rect(s13, 0, 0, W, H, MTU_BLUE)
# Gold strip
shape2 = s13.shapes.add_shape(1, Inches(0), Inches(4.7), W, Inches(0.06))
shape2.fill.solid(); shape2.fill.fore_color.rgb = MTU_GOLD
shape2.line.fill.background()

add_text(s13, "Thank You",
         Inches(0.6), Inches(1.2), Inches(12), Inches(1.4),
         font_size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s13, "Questions & Discussion",
         Inches(0.6), Inches(2.5), Inches(12), Inches(0.7),
         font_size=26, color=MTU_GOLD, align=PP_ALIGN.CENTER)

add_text(s13, "[Presenter Name]\n[Email]\nMunster Technological University",
         Inches(0.6), Inches(4.9), Inches(5.5), Inches(1.6),
         font_size=16, color=WHITE, align=PP_ALIGN.LEFT)

add_text(s13, "Repository: [GitHub / Balena link]\nDocs: ARCHITECTURE.md\nSmartCORE NWE: smartcore.nweurope.eu",
         Inches(6.8), Inches(4.9), Inches(6.0), Inches(1.6),
         font_size=14, color=RGBColor(0xBB, 0xCC, 0xDD), align=PP_ALIGN.LEFT)

add_rect(s13, Inches(0.6), Inches(6.7), Inches(2.6), Inches(0.65),
         RGBColor(0x00, 0x20, 0x60), line_color=MTU_GOLD, line_width=Pt(1.5))
add_text(s13, "[ MTU Logo ]", Inches(0.6), Inches(6.7), Inches(2.6), Inches(0.65),
         font_size=13, color=MTU_GOLD, align=PP_ALIGN.CENTER)
add_rect(s13, Inches(3.5), Inches(6.7), Inches(3.2), Inches(0.65),
         RGBColor(0x00, 0x20, 0x60), line_color=SC_GREEN, line_width=Pt(1.5))
add_text(s13, "[ SmartCORE NWE / Interreg Logo ]",
         Inches(3.5), Inches(6.7), Inches(3.2), Inches(0.65),
         font_size=11, color=SC_GREEN, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────────────────
out_path = "D:/ems/docs/SmartCORE_HEMS_Partner_Presentation.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Slides: {len(prs.slides)}")
