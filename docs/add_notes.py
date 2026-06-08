"""
Add detailed presenter notes to SmartCORE HEMS PowerPoint
Run after generate_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from lxml import etree

PPTX_PATH = "D:/ems/docs/SmartCORE_HEMS_Partner_Presentation.pptx"

# ── Notes text per slide (index 0 = slide 1) ─────────────────────────────────
NOTES = [

# ── SLIDE 1 — TITLE ───────────────────────────────────────────────────────────
"""Good morning / good afternoon everyone, and thank you for the opportunity to present today.

My name is [your name], I'm a researcher at Munster Technological University, and today I want to walk you through the Home Energy Management System — or HEMS box — that we've been developing here at MTU as part of the SmartCORE NWE project.

This is a ten-minute overview, so I'll move at a reasonable pace, but please do hold any questions until the end — I'll leave time for discussion.

What I want to show you today is a working prototype of a piece of hardware and software that sits inside a home, connects to the smart meter, monitors energy consumption in real time, and makes that data available for research and for the householder themselves.

Let's get started.""",

# ── SLIDE 2 — SMARTCORE CONTEXT ───────────────────────────────────────────────
"""Just to set the scene briefly for anyone who may be new to this specific workstream.

SmartCORE NWE is an Interreg North-West Europe project involving research partners from Ireland, the UK, Belgium, France, and the Netherlands. The overarching goal is to demonstrate community-scale smart energy management — helping households and communities reduce their energy costs and their carbon footprint.

Within that wider project, MTU leads the Ireland pilot workstream. Our specific job is to design and build the HEMS hardware and software stack — the physical box that goes into a home and does the energy monitoring.

We then target deploying these units into pilot homes across the Munster region. The numbers we're working towards are at least twenty homes, providing fifteen-minute resolution energy data that feeds into the SmartCORE European data platform.

One of the things I'll emphasise throughout this presentation is that everything we're building is open-source, remotely manageable, and designed so that any partner organisation could in principle pick up the same design and deploy it in their own country without starting from scratch.""",

# ── SLIDE 3 — THE PROBLEM ─────────────────────────────────────────────────────
"""So why are we building something ourselves rather than buying an off-the-shelf product? That's a fair question and I want to address it directly.

There are four core problems we identified.

First: the sensor ecosystem is fragmented. Smart home devices use different wireless protocols — Zigbee, Z-Wave, Bluetooth — and most commercial hubs don't give you raw data access. They give you a mobile app and a cloud subscription, and that's it.

Second: data sovereignty. For a research project, we need the raw energy data. We can't rely on a vendor's API that might change, or have data filtered or aggregated before we see it. And from a GDPR perspective, having home energy data flowing through a third-party cloud server in the US raises real compliance questions.

Third: remote manageability. Once we deploy twenty or thirty boxes across homes in Munster, we cannot send an engineer to each house every time there's a software update or a bug fix. The system has to update itself, heal itself, and keep running without any on-site intervention.

And fourth: we need an open, extensible research platform. Off-the-shelf products are black boxes. We need to be able to add new sensors, change data processing pipelines, export data in different formats, and integrate with the SmartCORE European platform — none of which is possible with consumer products.

Those four constraints are what drove us to build our own stack.""",

# ── SLIDE 4 — WHAT IS THE HEMS BOX ───────────────────────────────────────────
"""So what exactly is the HEMS box? Let me describe it in concrete terms.

On the hardware side, it's built around a Raspberry Pi 5. That's a small, low-power single-board computer — roughly the size of a credit card, runs on about five watts, and costs around eighty euros. It runs a specialised operating system called BalenaOS, which is purpose-built for running containerised applications on edge devices.

Plugged into the Pi via USB is a Sonoff ZBDongle-E — this is a Zigbee coordinator stick. Zigbee is the wireless protocol used by many smart home sensors, and this dongle acts as the radio gateway that lets the Pi talk to those sensors.

The sensor we're currently working with is the Frient SMSZB-120, which is a smart meter interface device. It clips onto the optical port of an Irish smart meter, counts the infrared blink pulses the meter emits, and converts those into real-time power readings in watts and cumulative energy readings in kilowatt-hours. It communicates back to the Pi over Zigbee.

On the software side, everything runs as Docker containers — eight separate services, all managed together using Docker Compose. I'll walk through what each of those services does in a moment.

The key design principles are: fully open-source throughout, no cloud dependency — the system works completely offline if needed — self-healing and self-updating via the BalenaOS infrastructure, extensible without hardware changes, and providing research-grade data access.

Think of it as a small, silent, self-managing computer that lives in a kitchen cupboard and watches the electricity meter.""",

# ── SLIDE 5 — SYSTEM ARCHITECTURE ────────────────────────────────────────────
"""Now let me show you the architecture — how all the pieces fit together.

Everything you see here runs on a single Raspberry Pi 5. The outer box represents the Pi itself, running BalenaOS. Inside that, we have a Docker bridge network — a private internal network with the IP range 172.18.4.0/24. Each container gets its own fixed IP address on that network, which makes the system predictable and easy to debug.

Starting from the bottom: the Sonoff ZBDongle-E USB stick is passed directly into the Home Assistant container. Next to it is the Frient sensor, communicating over Zigbee 802.15.4 radio.

In the top row, you have our four core services. Home Assistant at 172.18.4.2 is the central hub — it manages the Zigbee network, processes sensor data, and provides the user dashboard. InfluxDB at 172.18.4.3 is our time-series database where all historical energy data is stored. The Nginx reverse proxy at 172.18.4.4 sits in front of everything and provides a single secure entry point on port 80. And Mosquitto at 172.18.4.7 is our MQTT message broker — more on that in a moment.

In the second row: hass-configurator provides a browser-based editor for the Home Assistant configuration files. The led-status container drives three physical LEDs on the Pi's GPIO pins to give a visual health indication. The system-manager is a watchdog process that keeps everything running. And wifi-connect handles the initial WiFi setup when the box is first installed in a home.

On the right, you can see the BalenaOS Supervisor — the daemon that communicates with Balena Cloud and handles over-the-air updates.

The whole thing is entirely self-contained. Unplug it, move it to a different home, power it on, and it just works.""",

# ── SLIDE 6 — DATA FLOW ───────────────────────────────────────────────────────
"""Let me trace the path that energy data takes from the physical meter all the way through to the database, because I think this is where the elegance of the system becomes clear.

It starts with the smart meter itself. Irish smart meters emit an infrared pulse for every watt-hour of electricity consumed. The Frient sensor clips over that optical port and counts those pulses, converting them into power in watts and cumulative energy in kilowatt-hours.

The Frient sensor sends those readings wirelessly over Zigbee to the Sonoff USB dongle plugged into the Pi.

The dongle passes the Zigbee radio frames to Home Assistant via a USB serial connection, where the ZHA integration — that stands for Zigbee Home Automation — decodes them and updates the entity states in the Home Assistant state machine. So at this point, Home Assistant knows the current power reading.

The InfluxDB integration inside Home Assistant then automatically writes every state change to the InfluxDB database over HTTP, using InfluxDB's line protocol format. This happens without any custom code — it's built into Home Assistant.

And that's it. From physical pulse to queryable time-series data, with no custom middleware, no bespoke code, just configuration.

I want to flag one thing for the technically curious: you'll see at the bottom that we have a note about MQTT. We do have a Mosquitto MQTT broker running in the stack, but in the current ZHA deployment, it's actually idle — Zigbee data doesn't pass through it. We have an alternative architecture using Zigbee2MQTT that routes everything through MQTT, which opens up more advanced use cases, and I'll cover that on slide ten.""",

# ── SLIDE 7 — SERVICES DEEP DIVE ─────────────────────────────────────────────
"""Let me quickly walk through what each service actually does, because I want to give you a sense of how well-established and battle-tested each component is — we haven't had to write any of these from scratch.

Home Assistant is the largest open-source home automation project in the world, with over a million active installations. We use it as our data hub, Zigbee manager, and dashboard. It's the thing a householder would interact with.

InfluxDB is a purpose-built time-series database, widely used in industrial IoT and monitoring applications. Every sensor reading — timestamped to the second — goes in here and can be queried historically.

Nginx is a battle-hardened web server and reverse proxy. It's the gatekeeper — all web traffic comes in through it, it enforces basic authentication, and it routes requests to the right backend service.

The system-manager is the piece we wrote ourselves. It's a Python watchdog script that runs every thirty seconds. It checks whether a new software release has been pushed from Balena Cloud, and if so, gracefully restarts Home Assistant to pick up the changes. It also monitors internet connectivity and will reboot the Pi if the connection has been down for more than thirty minutes — because in a pilot home, we may not notice for days otherwise.

The led-status service is also ours. It drives three LEDs on the Pi's GPIO pins: green for Home Assistant healthy, blue for internet alive, and red if there are errors in the Home Assistant log. A green board is a happy board — even a non-technical householder can glance at it.

And Mosquitto is our MQTT message broker — lightweight, reliable, and the standard protocol for IoT messaging. It's ready and waiting for when we activate the Zigbee2MQTT path.""",

# ── SLIDE 8 — SELF-HEALING ────────────────────────────────────────────────────
"""One of the things I'm most proud of in this system is the self-healing capability, because it's what makes large-scale pilot deployment actually feasible.

Think about the logistics: twenty or thirty homes, spread across Munster, with a box running in each one. We cannot have a situation where one box gets a software bug and we have to drive out to the house to fix it. That's not sustainable.

BalenaOS solves the update problem completely. When we push a new Docker image to Balena Cloud, every device in the fleet automatically pulls the update and restarts the relevant container. Zero intervention required. We can fix a bug at eleven o'clock at night and have it deployed to all homes by midnight.

The system-manager watchdog handles the failure cases. If Home Assistant crashes or becomes unresponsive, the watchdog detects it and restarts it. If the internet goes down — which does happen in rural Irish homes — the watchdog waits thirty minutes and then reboots the entire Pi, which usually resolves connectivity issues caused by router or ISP problems.

The WiFi onboarding piece is important for the pilot deployment experience. When a box is first installed in a home, it has no WiFi credentials. What happens is: it broadcasts its own temporary WiFi network. The householder connects their phone to that network, gets a captive portal page — similar to what you'd see in a hotel — where they select their home WiFi and enter the password. The box connects, and that's the only time any technical setup is needed.

The physical LEDs are the last piece of the self-healing story. They give the householder a way to report back to us — "I can see the red light is on" — without needing to understand anything about the software.""",

# ── SLIDE 9 — DASHBOARD ───────────────────────────────────────────────────────
"""What does the householder actually see? This is the Home Assistant energy dashboard, accessible from any browser on their home network — or remotely once we have the mobile access solution in place.

[Point to the placeholder area] — I'll replace this placeholder with an actual screenshot before the final version of this deck, but what you'd see here is a real-time display of current power consumption in watts, a cumulative energy chart showing kilowatt-hours over time, and the ability to look back at historical data.

On the right you can see the data we're capturing: real-time power in watts, cumulative energy in kilowatt-hours, Zigbee signal quality so we know the sensor is communicating well, and all of this at sixty-second resolution stored in InfluxDB for the full history.

For research purposes, the full InfluxDB query interface is available at the /influx/ path, and we can query the data using the Flux query language, export it to CSV, or connect it directly to analysis tools. All of this is behind Nginx basic auth, so it's protected but accessible.""",

# ── SLIDE 10 — EXTENSIBILITY ──────────────────────────────────────────────────
"""Now I want to talk about where the system goes from here, because one of the key design decisions was to make this extensible without needing hardware changes.

The first extension I want to highlight is the Zigbee2MQTT alternative architecture. Currently, Zigbee data goes directly into Home Assistant via the ZHA integration. An alternative — which we've designed and documented but not yet activated — is to run a Zigbee2MQTT container instead. This routes all Zigbee data through the MQTT broker first. Why does that matter? Because with MQTT, multiple services can subscribe to the same data stream simultaneously. So you could have Home Assistant, a custom analytics service, and a cloud uploader all receiving the same sensor readings in parallel, without any of them interfering with the others. This is the architecture we'd want for a production deployment at scale.

For mobile access — so researchers or householders can check their energy data from outside the home — we have two options designed and ready to implement. The first is Tailscale, which creates a peer-to-peer VPN tunnel to the Pi. The second is Cloudflare Tunnel, which gives the box a public HTTPS URL without needing to open any ports on the home router. Both are activated by adding a single container and an API key — no code changes to the existing stack.

For future sensors: any Zigbee-compatible device pairs automatically. We're looking at temperature and humidity sensors, occupancy sensors, and integration with EV chargers. We also have a TCP proxy path in Nginx that's reserved for an eesmart D2L device — a more advanced smart meter interface — though the container for that hasn't been implemented yet.

And finally, for data export to the SmartCORE EU platform, the InfluxDB REST API is the natural integration point.""",

# ── SLIDE 11 — MTU CONTRIBUTION ───────────────────────────────────────────────
"""I want to be clear about what MTU has specifically contributed to this project, because I think it's significant.

We've built a complete, open-source, production-ready HEMS software stack from the ground up. Everything is containerised, reproducible, and version-controlled. Any engineer can clone the repository, run a single command, and have the entire system running — either on a Pi or on their own laptop for development.

We've validated this on Raspberry Pi 5 hardware specifically — and that's non-trivial, because the Pi 5 has architectural differences from previous Pi versions that required patches. In particular, the GPIO subsystem uses a different chip number, which broke the LED status software until we identified and fixed it.

We've written detailed technical documentation: a full architecture reference, a per-component guide, and a technical reference for the MQTT and mobile access architecture. This means that if any partner wants to understand the system, or if a new team member joins the project, the documentation is there.

In terms of research outputs, we're planning a peer-reviewed publication on the architecture, an anonymised dataset from the pilot homes, and a comparison study between the ZHA and Zigbee2MQTT approaches.

In terms of technology readiness — we're transitioning from TRL 4, which is lab validation, to TRL 5, which is pilot-validated. The core stack is functional and confirmed end-to-end on real hardware. Pilot deployment tooling is the next step.""",

# ── SLIDE 12 — STATUS & ROADMAP ───────────────────────────────────────────────
"""Let me bring you up to date on where we are right now and what the next twelve months look like.

Looking at the six status indicators across the top: the stack is running — all eight containers are operational on Pi 5 hardware. Sensor data is flowing end-to-end from the Frient sensor through to InfluxDB, confirmed in the lab. Over-the-air updates via BalenaOS have been validated. And the self-healing watchdog has been tested for both the HA restart and the device reboot scenarios.

Mobile access is in an intermediate state — the architecture is fully designed and documented, and the implementation is straightforward, but it's not yet activated. And pilot home deployment is where our focus is now — we're in the site survey phase and targeting first installations in Q3 of this year.

On the roadmap: Q2 2026 is about pilot preparation — finalising the deployment tooling, activating mobile access, and hardening the MQTT security before we put boxes into homes. Q3 is the first wave of pilot deployments, targeting five to ten homes with live data flowing to the SmartCORE platform. Q4 is about scaling to twenty homes and starting the energy baseline analysis. And by 2027 we're looking at full dataset analysis, a second-generation HEMS design based on what we learn, and sharing results across all NWE partner countries.

The short version: the system works, we're ready to deploy, and the focus now is execution.""",

# ── SLIDE 13 — THANK YOU ──────────────────────────────────────────────────────
"""And that brings me to the end of the presentation.

To summarise in three sentences: MTU has built an open-source, self-healing, containerised energy monitoring gateway on Raspberry Pi 5 hardware. It captures real-time energy data from smart meters via Zigbee sensors and stores it in a local time-series database. It's designed to be deployed in pilot homes with zero on-site maintenance and full remote manageability.

The system is working in the lab, and we're ready to start deploying.

I'm very happy to take questions now — whether on the technical architecture, the deployment approach, or how this integrates with what other partners are building. And if any of you want to go deeper on any aspect — particularly if you're thinking about replicating something similar in your own country — please do come and find me afterwards.

Thank you.""",

]

# ── Inject notes into slides ─────────────────────────────────────────────────
prs = Presentation(PPTX_PATH)

NSMAP = "http://schemas.openxmlformats.org/presentationml/2006/main"
ANSP  = "http://schemas.openxmlformats.org/drawingml/2006/main"

def set_notes(slide, text):
    """Add or replace the notes on a slide."""
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    # Clear existing paragraphs
    for p in tf.paragraphs:
        for r in p.runs:
            r.text = ""
    # Write new text line by line
    lines = text.strip().split("\n")
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.text = line

for i, (slide, note_text) in enumerate(zip(prs.slides, NOTES)):
    set_notes(slide, note_text)
    print(f"  Slide {i+1}: notes added ({len(note_text)} chars)")

OUT_PATH = PPTX_PATH.replace(".pptx", "_FINAL.pptx")
prs.save(OUT_PATH)
print(f"\nSaved with notes: {OUT_PATH}")
